import json
import logging
import re
from pathlib import Path

from app.config import BASE_DIR

logger = logging.getLogger("AICompanion.DocumentLoader")

KNOWLEDGE_DIR = BASE_DIR / "app" / "knowledge"
SOURCES_PATH = KNOWLEDGE_DIR / "sources.json"
INDEXES_DIR = KNOWLEDGE_DIR / "indexes"

SUPPORTED_EXTENSIONS = {".md", ".txt", ".pdf", ".docx", ".pptx"}


class DocumentLoader:
    def __init__(self):
        INDEXES_DIR.mkdir(parents=True, exist_ok=True)

    # ------------------------------------------------------------------ sources

    def load_active_sources(self) -> list:
        if not SOURCES_PATH.exists():
            logger.warning("sources.json not found.")
            return []
        try:
            with open(SOURCES_PATH, "r", encoding="utf-8") as f:
                sources = json.load(f)
        except json.JSONDecodeError as e:
            logger.error(f"sources.json is malformed: {e}")
            return []
        except Exception as e:
            logger.warning(f"Failed to read sources.json: {e}")
            return []
        if not isinstance(sources, list):
            logger.warning("sources.json must contain a list of source records.")
            return []
        return [s for s in sources if s.get("status") == "active"]

    def _active_path(self, source: dict) -> Path | None:
        active_version = source.get("active_version")
        for v in source.get("versions", []):
            if v.get("version") == active_version:
                return KNOWLEDGE_DIR / v["path"]
        return None

    # ------------------------------------------------------------------ extractors

    def _extract_text(self, file_path: Path) -> str:
        ext = file_path.suffix.lower()
        if ext in (".md", ".txt"):
            try:
                return file_path.read_text(encoding="utf-8")
            except (UnicodeDecodeError, OSError) as e:
                logger.warning(f"Failed to read {file_path.name}: {e}")
                return ""
        if ext == ".pdf":
            return self._extract_pdf(file_path)
        if ext == ".docx":
            return self._extract_docx(file_path)
        if ext == ".pptx":
            return self._extract_pptx(file_path)
        logger.warning(f"Unsupported file type: {ext} ({file_path.name})")
        return ""

    def _extract_pdf(self, file_path: Path) -> str:
        try:
            import pdfplumber
            pages = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append(f"[Halaman {i + 1}]\n{text.strip()}")
            return "\n\n".join(pages)
        except ImportError:
            logger.warning("pdfplumber not installed. Skipping PDF extraction.")
            return ""
        except Exception as e:
            logger.warning(f"Failed to extract PDF {file_path.name}: {e}")
            return ""

    def _extract_docx(self, file_path: Path) -> str:
        try:
            from docx import Document
            doc = Document(str(file_path))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            logger.warning("python-docx not installed. Skipping DOCX extraction.")
            return ""
        except Exception as e:
            logger.warning(f"Failed to extract DOCX {file_path.name}: {e}")
            return ""

    def _extract_pptx(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                texts.append(line)
                if texts:
                    slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except ImportError:
            logger.warning("python-pptx not installed. Skipping PPTX extraction.")
            return ""
        except Exception as e:
            logger.warning(f"Failed to extract PPTX {file_path.name}: {e}")
            return ""

    # ------------------------------------------------------------------ chunkers

    def _make_chunk(self, source: dict, index: int, heading: str, text: str) -> dict:
        return {
            "source_id": source.get("id"),
            "version": source.get("active_version"),
            "title": source.get("title"),
            "type": source.get("type"),
            "chunk_index": index,
            "heading": heading,
            "text": text,
        }

    def _chunk_markdown(self, text: str, source: dict) -> list:
        chunks = []
        current_heading = ""
        current_lines = []
        chunk_index = 0
        heading_found = False

        for line in text.splitlines():
            if re.match(r"^#{1,3}\s", line):
                heading_found = True
                content = "\n".join(current_lines).strip()
                if content:
                    chunks.append(self._make_chunk(source, chunk_index, current_heading, content))
                    chunk_index += 1
                current_heading = line.strip()
                current_lines = []
            else:
                current_lines.append(line)

        content = "\n".join(current_lines).strip()
        if content:
            chunks.append(self._make_chunk(source, chunk_index, current_heading, content))

        if not heading_found:
            return self._chunk_text(text, source)
        return chunks

    def _chunk_text(self, text: str, source: dict) -> list:
        chunks = []
        chunk_index = 0
        for para in re.split(r"\n{2,}", text):
            content = para.strip()
            if content:
                chunks.append(self._make_chunk(source, chunk_index, "", content))
                chunk_index += 1
        return chunks

    # ------------------------------------------------------------------ public

    def load_and_chunk(self, source: dict) -> list:
        file_path = self._active_path(source)
        if file_path is None:
            logger.warning(f"No active version path for source: {source.get('id')}")
            return []
        if not file_path.exists():
            logger.warning(f"File not found: {file_path}")
            return []

        ext = file_path.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            logger.warning(f"Skipping unsupported format: {ext} ({source.get('id')})")
            return []

        raw = self._extract_text(file_path)
        if not raw:
            return []

        if ext == ".md":
            return self._chunk_markdown(raw, source)
        return self._chunk_text(raw, source)

    def save_index(self, source_id: str, version: str, chunks: list) -> Path:
        index_path = INDEXES_DIR / f"{source_id}_{version}.json"
        try:
            with open(index_path, "w", encoding="utf-8") as f:
                json.dump(chunks, f, indent=2, ensure_ascii=False)
                f.write("\n")
            logger.info(f"Saved {len(chunks)} chunks -> {index_path.name}")
        except OSError as e:
            logger.warning(f"Could not save index cache for {source_id}: {e}")
        return index_path

    def load_index(self, source_id: str, version: str) -> list | None:
        index_path = INDEXES_DIR / f"{source_id}_{version}.json"
        if not index_path.exists():
            return None
        try:
            with open(index_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, OSError) as e:
            logger.warning(f"Corrupted index cache {index_path.name}: {e}. Will rechunk.")
            return None

    def load_all_documents(self, force_rechunk: bool = False) -> list:
        sources = self.load_active_sources()
        all_chunks = []

        for source in sources:
            source_id = source.get("id")
            version = source.get("active_version")

            if not force_rechunk:
                cached = self.load_index(source_id, version)
                if cached is not None:
                    logger.info(f"Using cached index: {source_id} {version} ({len(cached)} chunks)")
                    all_chunks.extend(cached)
                    continue

            chunks = self.load_and_chunk(source)
            if chunks:
                self.save_index(source_id, version, chunks)
                all_chunks.extend(chunks)
            else:
                logger.warning(f"No chunks produced for: {source_id}")

        return all_chunks
