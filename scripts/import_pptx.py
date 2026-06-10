"""
import_pptx.py — Impor file .pptx sebagai presentasi baru untuk Tenri.

Cara pakai:
    python scripts/import_pptx.py "C:/path/ke/file.pptx"
    python scripts/import_pptx.py "C:/path/ke/file.pptx" --name "Nama Presenter"

Yang dilakukan script ini:
    1. Baca semua slide dari file .pptx
    2. Salin file ke app/knowledge/presentations/
    3. Daftarkan ke app/knowledge/sources.json
    4. Buat app/presentation/slides.json baru
    5. Buat app/presentation/triggers.json baru
    6. Hapus cache index lama
    7. Cetak ringkasan hasil
"""

import argparse
import json
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

ROOT_DIR = Path(__file__).resolve().parent.parent
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

KNOWLEDGE_DIR   = ROOT_DIR / "app" / "knowledge"
SOURCES_PATH    = KNOWLEDGE_DIR / "sources.json"
INDEXES_DIR     = KNOWLEDGE_DIR / "indexes"
PPTX_DEST_DIR   = KNOWLEDGE_DIR / "presentations"
SLIDES_PATH     = ROOT_DIR / "app" / "presentation" / "slides.json"
TRIGGERS_PATH   = ROOT_DIR / "app" / "presentation" / "triggers.json"


# ------------------------------------------------------------------ helpers

def _slug(text: str) -> str:
    text = text.lower().strip()
    text = re.sub(r"[^\w\s-]", "", text)
    text = re.sub(r"[\s_]+", "-", text)
    return text[:40]


def _keywords(text: str, max_words: int = 5) -> list[str]:
    """Ambil kata bermakna dari teks untuk dijadikan pola trigger."""
    stop = {
        "dan", "atau", "yang", "di", "ke", "dari", "untuk", "pada", "dengan",
        "adalah", "ini", "itu", "juga", "dalam", "sebagai", "akan", "the",
        "and", "or", "of", "in", "to", "a", "an", "is", "for", "with",
    }
    words = re.findall(r"\b\w{3,}\b", text.lower())
    seen, result = set(), []
    for w in words:
        if w not in stop and w not in seen:
            seen.add(w)
            result.append(w)
        if len(result) >= max_words:
            break
    return result


def _read_pptx(pptx_path: Path) -> list[dict]:
    """Ekstrak data per slide dari file .pptx."""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body_lines = []
        notes_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if not line:
                        continue
                    # Heuristik: shape pertama dengan teks = judul
                    if not title and shape == slide.shapes[0]:
                        title = line
                    else:
                        body_lines.append(line)

        # Ambil catatan presenter jika ada
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                notes_text = notes_tf.text.strip()

        if not title and body_lines:
            title = body_lines.pop(0)

        slides.append({
            "slide_number": i,
            "title": title or f"Slide {i}",
            "body": body_lines,
            "notes": notes_text,
        })
    return slides


# ------------------------------------------------------------------ generators

def _make_slides_json(raw_slides: list[dict]) -> list[dict]:
    result = []
    for s in raw_slides:
        title = s["title"]
        body  = s["body"]
        notes = s["notes"]

        # Topics: kata kunci dari judul + baris pertama body
        topic_src = title + " " + " ".join(body[:3])
        topics = _keywords(topic_src, max_words=5)

        # Triggers untuk auto-advance: judul lowercase + keywords
        trig_words = list(dict.fromkeys(
            [title.lower()] + _keywords(title, 3) + _keywords(" ".join(body[:2]), 2)
        ))[:6]

        result.append({
            "id": s["slide_number"],
            "title": title,
            "topics": topics,
            "presenter_notes": notes or f"Jelaskan tentang {title}.",
            "tenri_role": "Berikan komentar atau keterangan singkat jika relevan.",
            "expected_tenri_mode": "comment",
            "triggers": trig_words,
        })
    return result


def _make_triggers_json(raw_slides: list[dict]) -> list[dict]:
    result = []
    modes_cycle = ["comment", "comment", "gentle_objection", "clarification",
                   "question", "memory", "grounding_check", "closing_memory"]

    for i, s in enumerate(raw_slides):
        title = s["title"]
        body  = s["body"]

        # Pola: judul utuh + keywords dari body
        patterns = list(dict.fromkeys(
            [title.lower()]
            + _keywords(title, 2)
            + _keywords(" ".join(body[:3]), 3)
        ))[:5]

        if not patterns:
            continue

        mode = modes_cycle[i % len(modes_cycle)]

        result.append({
            "id": f"t{i + 1:03d}",
            "patterns": patterns,
            "topic": title,
            "mode": mode,
            "slide_ids": [s["slide_number"]],
            "cooldown_seconds": 120,
            "suggested_response_intent": (
                f"Beri {mode.replace('_', ' ')} singkat terkait '{title}'. "
                "Maksimal 2 kalimat."
            ),
        })
    return result


def _load_sources() -> list:
    if not SOURCES_PATH.exists():
        return []
    try:
        with open(SOURCES_PATH, encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError):
        return []


def _save_sources(sources: list) -> None:
    with open(SOURCES_PATH, "w", encoding="utf-8") as f:
        json.dump(sources, f, indent=2, ensure_ascii=False)
        f.write("\n")


def _disable_all_sources() -> int:
    """Nonaktifkan semua sumber yang aktif. Kembalikan jumlah yang dinonaktifkan."""
    sources = _load_sources()
    count = 0
    for s in sources:
        if s.get("status") == "active":
            s["status"] = "disabled"
            count += 1
    _save_sources(sources)
    return count


def _update_sources(source_id: str, title: str, rel_path: str) -> None:
    sources = _load_sources()

    # Hapus entri lama dengan id yang sama jika ada
    sources = [s for s in sources if s.get("id") != source_id]

    sources.append({
        "id": source_id,
        "title": title,
        "type": "presentation",
        "active_version": "v1",
        "status": "active",
        "versions": [{
            "version": "v1",
            "path": rel_path,
            "status": "active",
            "created_at": datetime.now().strftime("%Y-%m-%d"),
        }],
        "notes": f"Diimpor otomatis oleh import_pptx.py pada {datetime.now().strftime('%Y-%m-%d %H:%M')}",
    })

    _save_sources(sources)


def _clear_indexes() -> int:
    if not INDEXES_DIR.exists():
        return 0
    deleted = 0
    for f in INDEXES_DIR.iterdir():
        if f.is_file() and f.suffix == ".json":
            f.unlink()
            deleted += 1
    return deleted


# ------------------------------------------------------------------ main

def run(pptx_path: Path, presenter_name: str = "", fresh: bool = False) -> None:
    if not pptx_path.exists():
        print(f"[ERROR] File tidak ditemukan: {pptx_path}")
        sys.exit(1)
    if pptx_path.suffix.lower() != ".pptx":
        print(f"[ERROR] File harus berformat .pptx, bukan {pptx_path.suffix}")
        sys.exit(1)

    print(f"\nMembaca: {pptx_path.name} ...")

    try:
        raw_slides = _read_pptx(pptx_path)
    except ImportError:
        print("[ERROR] python-pptx tidak terinstal. Jalankan: pip install python-pptx")
        sys.exit(1)
    except Exception as e:
        print(f"[ERROR] Gagal membaca file: {e}")
        sys.exit(1)

    if not raw_slides:
        print("[ERROR] Tidak ada slide yang berhasil dibaca dari file.")
        sys.exit(1)

    print(f"  {len(raw_slides)} slide ditemukan.")

    # 0. Mode pengetahuan
    if fresh:
        disabled = _disable_all_sources()
        if disabled:
            print(f"  Pengetahuan lama dinonaktifkan: {disabled} sumber")
        print("  Mode: hanya materi presentasi ini yang aktif (--fresh)")
    else:
        print("  Mode: semua pengetahuan tetap aktif — Tenri mengingat segalanya")

    # 1. Salin file ke knowledge/presentations/
    PPTX_DEST_DIR.mkdir(parents=True, exist_ok=True)
    dest_file = PPTX_DEST_DIR / pptx_path.name
    shutil.copy2(pptx_path, dest_file)
    rel_path = f"presentations/{pptx_path.name}"
    print(f"  File disalin → {dest_file.relative_to(ROOT_DIR)}")

    # 2. Daftarkan ke sources.json
    source_id = _slug(pptx_path.stem)
    title_label = (
        f"{pptx_path.stem} ({presenter_name})" if presenter_name else pptx_path.stem
    )
    _update_sources(source_id, title_label, rel_path)
    print(f"  Didaftarkan di sources.json dengan id: {source_id}")

    # 3. Buat slides.json
    slides_data = _make_slides_json(raw_slides)
    with open(SLIDES_PATH, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)
        f.write("\n")
    print(f"  slides.json → {len(slides_data)} slide")

    # 4. Buat triggers.json
    triggers_data = _make_triggers_json(raw_slides)
    with open(TRIGGERS_PATH, "w", encoding="utf-8") as f:
        json.dump(triggers_data, f, indent=2, ensure_ascii=False)
        f.write("\n")
    print(f"  triggers.json → {len(triggers_data)} trigger")

    # 5. Hapus cache index lama
    deleted = _clear_indexes()
    print(f"  Cache lama dihapus: {deleted} file")

    # 6. Ringkasan
    print("\n" + "=" * 55)
    print("  IMPOR SELESAI")
    print("=" * 55)
    print(f"  File       : {pptx_path.name}")
    print(f"  Slide      : {len(slides_data)}")
    print(f"  Trigger    : {len(triggers_data)}")
    print(f"  Mode       : {'fresh — hanya materi ini' if fresh else 'penuh — semua pengetahuan tetap aktif'}")
    print()
    print("  SLIDE YANG DIDETEKSI:")
    for s in slides_data:
        print(f"    {s['id']:2d}. {s['title']}")
    print()
    print("  LANGKAH SELANJUTNYA:")
    print("  1. Cek app/presentation/slides.json — periksa judul & topik")
    print("  2. Cek app/presentation/triggers.json — sesuaikan pola & mode")
    print("  3. Jalankan: python scripts/run_rehearsal.py")
    print("  4. Jalankan: python main.py")
    print("=" * 55 + "\n")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Impor file .pptx sebagai presentasi baru untuk Tenri.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Contoh:\n"
            "  python scripts/import_pptx.py presentasi.pptx\n"
            "  python scripts/import_pptx.py presentasi.pptx --name 'Budi'\n"
            "  python scripts/import_pptx.py presentasi.pptx --keep-existing\n"
        ),
    )
    parser.add_argument("pptx", help="Path ke file .pptx")
    parser.add_argument(
        "--name", default="",
        help="Nama presenter (opsional)"
    )
    parser.add_argument(
        "--fresh", action="store_true",
        help="Nonaktifkan semua materi lama, hanya gunakan presentasi ini"
    )
    args = parser.parse_args()
    run(Path(args.pptx), args.name, args.fresh)
