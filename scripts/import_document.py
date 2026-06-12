"""
import_document.py — Impor materi presentasi ke Tenri.

Format yang didukung: .pptx, .docx, .pdf
Cara pakai:
    python scripts/import_document.py "C:/path/ke/file.pptx"
    python scripts/import_document.py "C:/path/ke/file.pdf" --name "Nama Presenter"
    python scripts/import_document.py "C:/path/ke/file.docx" --fresh
"""

import argparse
import json
import logging
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

ROOT_DIR = Path(__file__).resolve().parent.parent
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

KNOWLEDGE_DIR  = ROOT_DIR / "app" / "knowledge"
SOURCES_PATH   = KNOWLEDGE_DIR / "sources.json"
INDEXES_DIR    = KNOWLEDGE_DIR / "indexes"
DEST_DIR       = KNOWLEDGE_DIR / "presentations"
SLIDES_PATH    = ROOT_DIR / "app" / "presentation" / "slides.json"
TRIGGERS_PATH  = ROOT_DIR / "app" / "presentation" / "triggers.json"

SUPPORTED = {".pptx", ".docx", ".pdf"}
logger = logging.getLogger("AICompanion.ImportDocument")


# ------------------------------------------------------------------ readers

def _read_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    sections = []
    for i, slide in enumerate(prs.slides, 1):
        title, body_lines, notes_text = "", [], ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if not line:
                        continue
                    if not title and shape == slide.shapes[0]:
                        title = line
                    else:
                        body_lines.append(line)
        if slide.has_notes_slide:
            nf = slide.notes_slide.notes_text_frame
            if nf:
                notes_text = nf.text.strip()
        if not title and body_lines:
            title = body_lines.pop(0)
        sections.append({
            "slide_number": i,
            "title": title or f"Slide {i}",
            "body": body_lines,
            "notes": notes_text,
        })
    return sections


def _is_docx_heading(para) -> bool:
    """Deteksi heading dari berbagai sinyal — tidak hanya dari nama style.

    Mengenali:
    1. Style Heading 1/2/3 atau Judul (cara standar Word)
    2. Seluruh paragraf di-bold (cara manual yang umum)
    3. Teks HURUF KAPITAL SEMUA yang pendek
    4. Ukuran font >= 14pt
    """
    text = para.text.strip()
    if not text or len(text) > 120:
        return False

    # 1. Nama style mengandung "heading", "judul", atau "title"
    style_name = (para.style.name or "").lower()
    if any(k in style_name for k in {"heading", "judul", "title", "head"}):
        return True

    # 2. Semua teks di-bold (abaikan run kosong)
    active_runs = [r for r in para.runs if r.text.strip()]
    if active_runs and all(r.bold for r in active_runs):
        return True

    # 3. Semua huruf kapital, panjang minimal 3 karakter, bukan angka semua
    if len(text) >= 3 and text == text.upper() and not text.replace(" ", "").isdigit():
        return True

    # 4. Ukuran font >= 14pt di salah satu run
    for run in active_runs:
        if run.font.size:
            try:
                from docx.shared import Pt
                if run.font.size >= Pt(14):
                    return True
            except Exception as exc:
                logger.warning(
                    "Gagal memeriksa ukuran font heading DOCX untuk paragraf %r: %s",
                    text[:60],
                    exc,
                )

    return False


def _read_docx(path: Path) -> list[dict]:
    from docx import Document
    doc = Document(str(path))

    sections, current_title, current_body = [], "", []
    section_num = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if _is_docx_heading(para):
            if current_title or current_body:
                section_num += 1
                sections.append({
                    "slide_number": section_num,
                    "title": current_title or f"Bagian {section_num}",
                    "body": current_body,
                    "notes": "",
                })
            current_title = text
            current_body = []
        else:
            current_body.append(text)

    # Simpan section terakhir
    if current_title or current_body:
        section_num += 1
        sections.append({
            "slide_number": section_num,
            "title": current_title or f"Bagian {section_num}",
            "body": current_body,
            "notes": "",
        })

    # Fallback: tidak ada heading terdeteksi — pecah setiap 5 paragraf
    if not sections:
        all_paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        chunk_size = 5
        for i in range(0, len(all_paras), chunk_size):
            chunk = all_paras[i:i + chunk_size]
            num = i // chunk_size + 1
            sections.append({
                "slide_number": num,
                "title": chunk[0] if chunk else f"Bagian {num}",
                "body": chunk[1:],
                "notes": "",
            })

    return sections


def _read_pdf(path: Path) -> list[dict]:
    import pdfplumber
    sections = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            raw = page.extract_text()
            if not raw or not raw.strip():
                continue
            lines = [l.strip() for l in raw.splitlines() if l.strip()]
            if not lines:
                continue
            # Baris pertama yang cukup panjang = judul kandidat
            title = ""
            body_start = 0
            for j, line in enumerate(lines):
                if len(line) >= 5:
                    title = line
                    body_start = j + 1
                    break
            body = lines[body_start:]
            sections.append({
                "slide_number": i,
                "title": title or f"Halaman {i}",
                "body": body,
                "notes": "",
            })
    return sections


def _read_document(path: Path) -> list[dict]:
    ext = path.suffix.lower()
    try:
        if ext == ".pptx":
            return _read_pptx(path)
        if ext == ".docx":
            return _read_docx(path)
        if ext == ".pdf":
            return _read_pdf(path)
    except ImportError as e:
        missing = str(e).split("'")[1] if "'" in str(e) else str(e)
        lib_map = {
            "pptx": "python-pptx",
            "docx": "python-docx",
            "pdfplumber": "pdfplumber",
        }
        lib = next((v for k, v in lib_map.items() if k in missing), missing)
        raise ModuleNotFoundError(f"Library tidak terinstal: pip install {lib}") from e
    raise ValueError(f"Format tidak didukung: {ext}")


# ------------------------------------------------------------------ generators

def _slug(text: str) -> str:
    text = text.lower().strip()
    text = re.sub(r"[^\w\s-]", "", text)
    text = re.sub(r"[\s_]+", "-", text)
    return text[:40]


def _keywords(text: str, max_words: int = 5) -> list[str]:
    stop = {
        "dan", "atau", "yang", "di", "ke", "dari", "untuk", "pada", "dengan",
        "adalah", "ini", "itu", "juga", "dalam", "sebagai", "akan", "the",
        "and", "or", "of", "in", "to", "a", "an", "is", "for", "with",
        "pada", "agar", "itu", "ini", "saat", "setelah", "karena",
    }
    words = re.findall(r"\b\w{3,}\b", text.lower())
    seen, result = set(), []
    for w in words:
        if w not in stop and w not in seen:
            seen.add(w)
            result.append(w)
        if len(result) >= max_words:
            break
    return result


_GENERIC_SINGLE_WORD_TRIGGERS = frozenset({
    "awal", "akhir", "bagian", "dasar", "definisi", "kecerdasan",
    "generatif", "istilah", "kesimpulan", "kritis", "manfaat", "mengenal",
    "pelatihan", "pemetaan", "pengenalan", "penutup", "ringkasan", "risiko",
    "sikap", "teknologi", "tepat", "umum",
})


def _safe_trigger_patterns(patterns: list[str], limit: int = 5) -> list[str]:
    """Keep specific phrases and domain terms; reject generic one-word triggers."""
    result: list[str] = []
    seen: set[str] = set()
    for pattern in patterns:
        normalized = re.sub(r"\s+", " ", pattern.lower()).strip()
        if not normalized or normalized in seen:
            continue
        if " " not in normalized and normalized in _GENERIC_SINGLE_WORD_TRIGGERS:
            continue
        seen.add(normalized)
        result.append(normalized)
        if len(result) >= limit:
            break
    return result


def _make_slides_json(sections: list[dict]) -> list[dict]:
    result = []
    for s in sections:
        title = s["title"]
        body  = s["body"]
        notes = s["notes"]
        topic_src = title + " " + " ".join(body[:3])
        topics = _keywords(topic_src, 5)
        trig_words = list(dict.fromkeys(
            [title.lower()] + _keywords(title, 3) + _keywords(" ".join(body[:2]), 2)
        ))[:6]
        result.append({
            "id": s["slide_number"],
            "title": title,
            "topics": topics,
            "presenter_notes": notes or f"Jelaskan tentang {title}.",
            "tenri_role": "Berikan komentar atau keterangan singkat jika relevan.",
            "expected_tenri_mode": "comment",
            "triggers": trig_words,
        })
    return result


def _make_triggers_json(sections: list[dict]) -> list[dict]:
    modes_cycle = ["comment", "comment", "gentle_objection", "clarification",
                   "question", "memory", "grounding_check", "closing_memory"]
    result = []
    for i, s in enumerate(sections):
        title, body = s["title"], s["body"]
        patterns = _safe_trigger_patterns(
            [title.lower()]
            + _keywords(title, 2)
            + _keywords(" ".join(body[:3]), 3)
        )
        if not patterns:
            continue
        mode = modes_cycle[i % len(modes_cycle)]
        result.append({
            "id": f"t{i + 1:03d}",
            "patterns": patterns,
            "topic": title,
            "mode": mode,
            "slide_ids": [s["slide_number"]],
            "cooldown_seconds": 120,
            "suggested_response_intent": (
                f"Beri {mode.replace('_', ' ')} singkat terkait '{title}'. "
                "Maksimal 2 kalimat."
            ),
        })
    return result


# ------------------------------------------------------------------ sources helpers

def _load_sources() -> list:
    if not SOURCES_PATH.exists():
        return []
    try:
        with open(SOURCES_PATH, encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError):
        return []


def _save_sources(sources: list) -> None:
    with open(SOURCES_PATH, "w", encoding="utf-8") as f:
        json.dump(sources, f, indent=2, ensure_ascii=False)
        f.write("\n")


def _disable_all_sources() -> int:
    sources = _load_sources()
    count = sum(1 for s in sources if s.get("status") == "active")
    for s in sources:
        if s.get("status") == "active":
            s["status"] = "disabled"
    _save_sources(sources)
    return count


def _register_source(source_id: str, title: str, rel_path: str,
                     file_type: str = "presentation") -> None:
    sources = _load_sources()
    sources = [s for s in sources if s.get("id") != source_id]
    sources.append({
        "id": source_id,
        "title": title,
        "type": file_type,
        "active_version": "v1",
        "status": "active",
        "versions": [{
            "version": "v1",
            "path": rel_path,
            "status": "active",
            "created_at": datetime.now().strftime("%Y-%m-%d"),
        }],
        "notes": (
            f"Diimpor otomatis oleh import_document.py "
            f"pada {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        ),
    })
    _save_sources(sources)


def _clear_indexes() -> int:
    if not INDEXES_DIR.exists():
        return 0
    deleted = 0
    for f in INDEXES_DIR.iterdir():
        if f.is_file() and f.suffix == ".json":
            f.unlink()
            deleted += 1
    return deleted


# ------------------------------------------------------------------ main

def run(doc_path: Path, presenter_name: str = "", fresh: bool = False) -> None:
    ext = doc_path.suffix.lower()

    # Raise (bukan sys.exit) supaya pemanggil dari menu bisa menangkapnya dan
    # kembali ke menu. Blok __main__ CLI menerjemahkan exception ke exit code.
    if not doc_path.exists():
        raise FileNotFoundError(f"File tidak ditemukan: {doc_path}")
    if ext not in SUPPORTED:
        raise ValueError(f"Format tidak didukung: {ext}. Gunakan .pptx, .docx, atau .pdf")

    print(f"\nMembaca: {doc_path.name} ({ext.upper()[1:]}) ...")

    sections = _read_document(doc_path)

    if not sections:
        raise ValueError("Tidak ada konten yang berhasil dibaca dari file.")

    print(f"  {len(sections)} bagian/halaman ditemukan.")

    # Nonaktifkan pengetahuan lama jika mode fresh
    if fresh:
        disabled = _disable_all_sources()
        if disabled:
            print(f"  Pengetahuan lama dinonaktifkan: {disabled} sumber")
        print("  Mode: hanya materi ini yang aktif (--fresh)")
    else:
        print("  Mode: semua pengetahuan tetap aktif — Tenri mengingat segalanya")

    # Salin file
    DEST_DIR.mkdir(parents=True, exist_ok=True)
    dest_file = DEST_DIR / doc_path.name
    shutil.copy2(doc_path, dest_file)
    rel_path = f"presentations/{doc_path.name}"
    print(f"  File disalin → {dest_file.relative_to(ROOT_DIR)}")

    # Daftarkan ke sources.json
    source_id = _slug(doc_path.stem)
    title_label = (
        f"{doc_path.stem} ({presenter_name})" if presenter_name else doc_path.stem
    )
    _register_source(source_id, title_label, rel_path)
    print(f"  Didaftarkan di sources.json dengan id: {source_id}")

    # Buat slides.json
    slides_data = _make_slides_json(sections)
    with open(SLIDES_PATH, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)
        f.write("\n")
    print(f"  slides.json → {len(slides_data)} bagian")

    # Buat triggers.json
    triggers_data = _make_triggers_json(sections)
    with open(TRIGGERS_PATH, "w", encoding="utf-8") as f:
        json.dump(triggers_data, f, indent=2, ensure_ascii=False)
        f.write("\n")
    print(f"  triggers.json → {len(triggers_data)} trigger")

    # Hapus cache index lama
    deleted = _clear_indexes()
    print(f"  Cache lama dihapus: {deleted} file")

    # Ringkasan
    print("\n" + "=" * 55)
    print("  IMPOR SELESAI")
    print("=" * 55)
    print(f"  File   : {doc_path.name}")
    print(f"  Format : {ext.upper()[1:]}")
    print(f"  Bagian : {len(slides_data)}")
    print(f"  Mode   : {'fresh' if fresh else 'penuh — semua pengetahuan aktif'}")
    print()
    print("  BAGIAN YANG DIDETEKSI:")
    for s in slides_data[:10]:
        print(f"    {s['id']:2d}. {s['title']}")
    if len(slides_data) > 10:
        print(f"    ... dan {len(slides_data) - 10} bagian lainnya")
    print()
    print("  LANGKAH SELANJUTNYA:")
    print("  1. Cek app/presentation/slides.json — periksa judul & topik")
    print("  2. Cek app/presentation/triggers.json — sesuaikan pola & mode")
    print("  3. Jalankan: python scripts/run_rehearsal.py")
    print("  4. Jalankan: python main.py")
    print("=" * 55 + "\n")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Impor materi presentasi (.pptx/.docx/.pdf) ke Tenri.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Contoh:\n"
            "  python scripts/import_document.py presentasi.pptx\n"
            "  python scripts/import_document.py makalah.pdf --name 'Budi'\n"
            "  python scripts/import_document.py laporan.docx --fresh\n"
        ),
    )
    parser.add_argument("file", help="Path ke file (.pptx, .docx, atau .pdf)")
    parser.add_argument("--name", default="", help="Nama presenter (opsional)")
    parser.add_argument(
        "--fresh", action="store_true",
        help="Nonaktifkan materi lama, hanya gunakan file ini"
    )
    args = parser.parse_args()
    try:
        run(Path(args.file), args.name, args.fresh)
    except Exception as e:
        print(f"[ERROR] {e}")
        sys.exit(1)
